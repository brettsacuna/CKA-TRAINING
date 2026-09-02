"""
Sistema de diseño para las diapositivas del curso CKA Hands-On Training.

Tema: oscuro, "Kubernetes blue". Todo el estilo se aplica de forma explícita
sobre cada shape (no se depende del tema de Office) para que el render sea
idéntico en cualquier equipo.

Unidades: EMU (914400 por pulgada). Lienzo 16:9 -> 12192000 x 6858000.
"""

from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --------------------------------------------------------------------------- #
# Lienzo y rejilla
# --------------------------------------------------------------------------- #
EMU_W = 12192000
EMU_H = 6858000
MARGIN = 620000                      # margen izquierdo / derecho
CONTENT_W = EMU_W - 2 * MARGIN       # ancho útil
COL_GAP = 320000                     # separación entre columnas

Y_EYEBROW = 300000
Y_HEADING = 560000
Y_RULE = 1230000                     # regla de acento bajo el título
Y_BODY = 1470000                     # inicio del cuerpo
Y_FOOTER = 6430000

# --------------------------------------------------------------------------- #
# Paleta
# --------------------------------------------------------------------------- #
BG          = RGBColor(0x0B, 0x12, 0x20)   # fondo base (casi negro azulado)
BG_ALT      = RGBColor(0x0E, 0x17, 0x2A)   # fondo de secciones / divisores
PANEL       = RGBColor(0x14, 0x1E, 0x33)   # tarjetas y paneles
PANEL_SOFT  = RGBColor(0x11, 0x1A, 0x2C)   # paneles secundarios
HAIRLINE    = RGBColor(0x26, 0x35, 0x4E)   # bordes y separadores

INK         = RGBColor(0xF3, 0xF6, 0xFC)   # títulos
BODY        = RGBColor(0xBD, 0xC9, 0xDB)   # texto corrido
MUTED       = RGBColor(0x84, 0x93, 0xAD)   # pies, metadatos
FAINT       = RGBColor(0x4C, 0x5B, 0x75)   # números fantasma, motivos

K8S         = RGBColor(0x32, 0x6C, 0xE5)   # azul Kubernetes (rellenos)
K8S_TEXT    = RGBColor(0x74, 0xA6, 0xFF)   # azul legible sobre fondo oscuro
K8S_DEEP    = RGBColor(0x1E, 0x40, 0x8F)   # azul profundo (bordes de relleno)

OK          = RGBColor(0x3F, 0xB9, 0x50)
OK_TEXT     = RGBColor(0x63, 0xD1, 0x78)
WARN        = RGBColor(0xE3, 0xB3, 0x41)
WARN_TEXT   = RGBColor(0xF0, 0xC9, 0x6B)
DANGER      = RGBColor(0xF8, 0x51, 0x49)
DANGER_TEXT = RGBColor(0xFF, 0x7B, 0x72)

CODE_BG      = RGBColor(0x0A, 0x11, 0x20)
CODE_TEXT    = RGBColor(0xE6, 0xED, 0xF3)
CODE_COMMENT = RGBColor(0x7C, 0x8A, 0xA5)
CODE_ACCENT  = RGBColor(0x9C, 0xDC, 0xFE)   # tokens destacados (prompts $, k)

# Acentos disponibles por nombre (para variar por sección)
ACCENTS = {
    "blue":   (K8S, K8S_TEXT, K8S_DEEP),
    "ok":     (OK, OK_TEXT, RGBColor(0x1F, 0x6F, 0x2B)),
    "warn":   (WARN, WARN_TEXT, RGBColor(0x8A, 0x67, 0x12)),
    "danger": (DANGER, DANGER_TEXT, RGBColor(0x8E, 0x2B, 0x25)),
}

# --------------------------------------------------------------------------- #
# Tipografía
# --------------------------------------------------------------------------- #
F_HEAD   = "Segoe UI Semibold"
F_LIGHT  = "Segoe UI Light"
F_BODY   = "Segoe UI"
F_MONO   = "Consolas"

# Escala de tamaños (pt)
SZ_EYEBROW = 11
SZ_H1      = 30      # título de diapositiva de contenido
SZ_H_SECTION = 54    # título de divisor de sección
SZ_LEAD    = 15
SZ_BODY    = 13
SZ_SMALL   = 11
SZ_CODE    = 12
SZ_FOOTER  = 9


# --------------------------------------------------------------------------- #
# Utilidades de bajo nivel
# --------------------------------------------------------------------------- #
def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _no_autofit(tf):
    # desactiva el autoajuste para que el tamaño de fuente sea fiable
    el = tf._txBody
    bodyPr = el.find(qn("a:bodyPr"))
    for tag in ("a:normAutofit", "a:spAutoFit"):
        e = bodyPr.find(qn(tag))
        if e is not None:
            bodyPr.remove(e)
    bodyPr.append(el.makeelement(qn("a:noAutofit"), {}))


def _spacing(run, pts):
    """Interletrado (tracking) en puntos."""
    run._r.get_or_add_rPr().set("spc", str(int(pts * 100)))


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = box.text_frame
    tf.word_wrap = wrap
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box


def _style_run(r, *, size, color, font=F_BODY, bold=False, italic=False,
               spacing=None):
    r.font.size = Pt(size)
    r.font.name = font
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if spacing is not None:
        _spacing(r, spacing)


def para(tf, text, *, size=SZ_BODY, color=BODY, font=F_BODY, bold=False,
         italic=False, align=PP_ALIGN.LEFT, space_after=6, space_before=0,
         line=1.18, spacing=None, first=False, bullet=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    if space_after is not None:
        p.space_after = Pt(space_after)
    if space_before is not None:
        p.space_before = Pt(space_before)
    if line is not None:
        p.line_spacing = line
    _set_bullet(p, bullet)
    r = p.add_run()
    r.text = text
    _style_run(r, size=size, color=color, font=font, bold=bold, italic=italic,
               spacing=spacing)
    return p


def _set_bullet(p, on):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone", "a:buFont"):
        e = pPr.find(qn(tag))
        if e is not None:
            pPr.remove(e)
    if on:
        pPr.set("marL", "182880")
        pPr.set("indent", "-182880")
        buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "•"})
        pPr.append(buFont)
        pPr.append(buChar)
    else:
        pPr.set("marL", "0")
        pPr.set("indent", "0")
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def rect(slide, x, y, w, h, *, fill=PANEL, line=None, line_w=1.0,
         rounded=False, radius=0.045, shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    _no_autofit(shp.text_frame)
    for m in ("left", "right", "top", "bottom"):
        setattr(shp.text_frame, f"margin_{m}", Emu(0))
    return shp


def oval(slide, x, y, w, h, *, fill=K8S, line=None, line_w=1.25):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(x)), Emu(int(y)),
                                 Emu(int(w)), Emu(int(h)))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def chip(slide, x, y, w, h, text, *, fill=None, line=HAIRLINE, text_color=BODY,
         size=SZ_SMALL, bold=True, font=F_BODY, align=PP_ALIGN.CENTER,
         spacing=0.4):
    shp = rect(slide, x, y, w, h, fill=fill, line=line, line_w=1.0,
               rounded=True, radius=0.5)
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, text, size=size, color=text_color, font=font, bold=bold,
         align=align, first=True, space_after=0, line=1.0, spacing=spacing)
    return shp


def num_badge(slide, x, y, d, label, *, fill=K8S, text_color=INK, size=12):
    o = oval(slide, x, y, d, d, fill=fill)
    tf = o.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _no_autofit(tf)
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", Emu(0))
    para(tf, str(label), size=size, color=text_color, font=F_HEAD, bold=True,
         align=PP_ALIGN.CENTER, first=True, space_after=0, line=1.0)
    return o


# --------------------------------------------------------------------------- #
# Componentes de página
# --------------------------------------------------------------------------- #
def eyebrow(slide, text, *, accent_text=K8S_TEXT, x=MARGIN, w=CONTENT_W):
    box = textbox(slide, x, Y_EYEBROW, w, 240000, anchor=MSO_ANCHOR.MIDDLE)
    para(box.text_frame, text.upper(), size=SZ_EYEBROW, color=accent_text,
         font=F_HEAD, bold=True, spacing=2.4, first=True, space_after=0, line=1.0)
    return box


def heading(slide, text, *, x=MARGIN, w=CONTENT_W, size=SZ_H1):
    box = textbox(slide, x, Y_HEADING, w, 620000, anchor=MSO_ANCHOR.TOP)
    para(box.text_frame, text, size=size, color=INK, font=F_LIGHT, bold=False,
         first=True, space_after=0, line=1.05)
    return box


def accent_rule(slide, *, x=MARGIN, y=Y_RULE, w=560000, color=K8S):
    rect(slide, x, y, w, 46000, fill=color, line=None)


def footer(slide, clase_n, page):
    rect(slide, MARGIN, Y_FOOTER - 26000, CONTENT_W, 12000, fill=HAIRLINE, line=None)
    left = textbox(slide, MARGIN, Y_FOOTER, CONTENT_W * 0.7, 220000,
                   anchor=MSO_ANCHOR.MIDDLE)
    para(left.text_frame, f"CKA Hands-On Training  ·  Clase {clase_n}",
         size=SZ_FOOTER, color=MUTED, font=F_BODY, first=True, space_after=0,
         line=1.0, spacing=0.3)
    right = textbox(slide, MARGIN + CONTENT_W * 0.7, Y_FOOTER, CONTENT_W * 0.3,
                    220000, anchor=MSO_ANCHOR.MIDDLE)
    para(right.text_frame, f"{page:02d}", size=SZ_FOOTER, color=MUTED,
         font=F_HEAD, bold=True, align=PP_ALIGN.RIGHT, first=True,
         space_after=0, line=1.0, spacing=0.6)


def hex_motif(slide, *, cx, cy, r, color=FAINT, line_w=1.5):
    """Heptágono tipo timón de Kubernetes, muy sutil, como marca de agua."""
    import math
    pts = []
    for i in range(7):
        a = -math.pi / 2 + i * (2 * math.pi / 7)
        pts.append((cx + r * math.cos(a), cy + r * math.sin(a)))
    free = slide.shapes.build_freeform(Emu(int(pts[0][0])), Emu(int(pts[0][1])))
    free.add_line_segments([(Emu(int(x)), Emu(int(y))) for x, y in pts[1:]],
                           close=True)
    shp = free.convert_to_shape()
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def code_panel(slide, x, y, w, lines, *, accent=K8S, min_h=None, pad=190000,
               size=SZ_CODE, title=None):
    """Panel de código. `lines` es lista de str. Comentarios (#...) atenuados."""
    line_h = int(Pt(size).emu * 1.42)
    head_h = 300000 if title else 0
    body_h = int(line_h * max(1, len(lines))) + 2 * pad
    h = body_h + head_h
    if min_h:
        h = max(h, min_h)
    panel = rect(slide, x, y, w, h, fill=CODE_BG, line=HAIRLINE, line_w=1.0,
                 rounded=True, radius=0.03)
    rect(slide, x, y + (head_h if title else 0), 46000,
         h - (head_h if title else 0), fill=accent, line=None)
    if title:
        tb = textbox(slide, x + pad, y + 60000, w - 2 * pad, head_h - 40000,
                     anchor=MSO_ANCHOR.MIDDLE)
        para(tb.text_frame, title.upper(), size=SZ_SMALL - 1, color=MUTED,
             font=F_HEAD, bold=True, spacing=1.8, first=True, space_after=0,
             line=1.0)
        rect(slide, x + pad, y + head_h - 6000, w - 2 * pad, 10000,
             fill=HAIRLINE, line=None)
    tb = textbox(slide, x + pad + 40000, y + head_h + pad, w - 2 * pad - 40000,
                 body_h - 2 * pad)
    tf = tb.text_frame
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.34
        p.space_after = Pt(0)
        _set_bullet(p, False)
        _emit_code_line(p, ln, size)
    return panel


def _emit_code_line(p, ln, size):
    """Divide una línea de código en código + comentario para colorearlos."""
    if ln.strip().startswith("#"):
        r = p.add_run(); r.text = ln
        _style_run(r, size=size, color=CODE_COMMENT, font=F_MONO, italic=True)
        return
    idx = ln.find("  #")
    if idx == -1 and " #" in ln and "#" not in ln.split(" #")[0]:
        idx = ln.find(" #")
    if idx != -1:
        code, comment = ln[:idx], ln[idx:]
        r = p.add_run(); r.text = code
        _style_run(r, size=size, color=CODE_TEXT, font=F_MONO)
        r = p.add_run(); r.text = comment
        _style_run(r, size=size, color=CODE_COMMENT, font=F_MONO, italic=True)
    else:
        r = p.add_run(); r.text = ln
        _style_run(r, size=size, color=CODE_TEXT, font=F_MONO)


def note_band(slide, text, *, y, accent=K8S, x=MARGIN, w=CONTENT_W, h=470000,
              label=None, label_color=K8S_TEXT):
    rect(slide, x, y, w, h, fill=PANEL_SOFT, line=HAIRLINE, line_w=1.0,
         rounded=True, radius=0.10)
    rect(slide, x + 30000, y + 24000, 40000, h - 48000, fill=accent, line=None)
    tb = textbox(slide, x + 240000, y, w - 470000, h, anchor=MSO_ANCHOR.MIDDLE)
    tf = tb.text_frame
    if label:
        para(tf, label.upper(), size=SZ_SMALL - 1, color=label_color,
             font=F_HEAD, bold=True, spacing=1.8, first=True, space_after=4,
             line=1.0)
    para(tf, text, size=SZ_SMALL, color=BODY, font=F_BODY, italic=True,
         first=not label, space_after=0, line=1.16)


def scaffold(slide, *, clase_n, page, eyebrow_text, heading_text,
             accent_text=K8S_TEXT, rule_color=K8S, motif=False, heading_w=None):
    """Estructura común: fondo, eyebrow, título, regla y pie."""
    set_bg(slide)
    eyebrow(slide, eyebrow_text, accent_text=accent_text)
    heading(slide, heading_text, w=heading_w or CONTENT_W)
    accent_rule(slide, color=rule_color)
    footer(slide, clase_n, page)
