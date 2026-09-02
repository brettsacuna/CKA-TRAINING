# -*- coding: utf-8 -*-
"""
Iconos y grafismos vectoriales para las diapositivas.

Todo se dibuja con formas de PowerPoint (líneas, óvalos, freeforms) — sin
archivos de imagen — para que quede nítido y en el tema oscuro del curso.

Convención de cada icono:  draw(slide, x, y, s, color, w=None)
    (x, y) esquina superior izquierda del cuadro,  s = lado,  w = grosor de trazo.
"""

import math
import os
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

import design as D

# --------------------------------------------------------------------------- #
# Iconos oficiales de Kubernetes  (PNG a color)
# Fuente: github.com/kubernetes/community/icons  ·  Apache-2.0 / CC-BY-4.0
# Ver assets/k8s/NOTICE.md
# --------------------------------------------------------------------------- #
_ASSET = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "k8s")

_PNG = {
    "pod": "pod.png", "deploy": "deploy.png", "rs": "rs.png", "sts": "sts.png",
    "ds": "ds.png", "svc": "svc.png", "ep": "ep.png", "ing": "ing.png",
    "netpol": "netpol.png", "pv": "pv.png", "pvc": "pvc.png", "vol": "vol.png",
    "sc": "sc.png", "cm": "cm.png", "secret": "secret.png", "role": "role.png",
    "crole": "c-role.png", "sa": "sa.png", "crd": "crd.png", "ns": "ns.png",
    "hpa": "hpa.png", "limits": "limits.png", "quota": "quota.png",
    "job": "job.png", "cronjob": "cronjob.png",
    "etcd": "etcd.png", "node": "node.png", "controlplane": "control-plane.png",
    "master": "master.png", "sched": "cp-sched.png", "apiserver": "cp-api.png",
    "kubelet": "cp-kubelet.png", "kproxy": "cp-k-proxy.png",
    "wheel": "kubernetes.png",
}


def _png_path(slot):
    f = _PNG.get(slot)
    if not f:
        return None
    p = os.path.join(_ASSET, f)
    return p if os.path.exists(p) else None


# --------------------------------------------------------------------------- #
# Primitivas de trazo
# --------------------------------------------------------------------------- #
def _bare(shape, color, w, fill=None):
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = color
    shape.line.width = Pt(w)
    shape.shadow.inherit = False
    return shape


def ring(slide, cx, cy, r, color, w):
    return _bare(slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Emu(int(cx - r)), Emu(int(cy - r)),
        Emu(int(2 * r)), Emu(int(2 * r))), color, w)


def dot(slide, cx, cy, r, color):
    return _bare(slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Emu(int(cx - r)), Emu(int(cy - r)),
        Emu(int(2 * r)), Emu(int(2 * r))), color, 0.5, fill=color)


def rrect(slide, x, y, w_, h_, color, w, fill=None, rad=0.18):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(x)),
                                Emu(int(y)), Emu(int(w_)), Emu(int(h_)))
    try:
        sh.adjustments[0] = rad
    except Exception:
        pass
    return _bare(sh, color, w, fill=fill)


def line(slide, x1, y1, x2, y2, color, w, arrow=False, dashed=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(x1)),
                                    Emu(int(y1)), Emu(int(x2)), Emu(int(y2)))
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    cn.shadow.inherit = False
    ln = cn.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    if arrow:
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))
    return cn


def poly(slide, pts, color, w, fill=None, close=True):
    fb = slide.shapes.build_freeform(Emu(int(pts[0][0])), Emu(int(pts[0][1])),
                                     scale=1)
    fb.add_line_segments([(Emu(int(px)), Emu(int(py))) for px, py in pts[1:]],
                         close=close)
    return _bare(fb.convert_to_shape(), color, w, fill=fill)


def arc(slide, cx, cy, r, a0, a1, color, w):
    sh = slide.shapes.add_shape(MSO_SHAPE.ARC, Emu(int(cx - r)), Emu(int(cy - r)),
                                Emu(int(2 * r)), Emu(int(2 * r)))
    try:
        sh.adjustments[0] = a0
        sh.adjustments[1] = a1
    except Exception:
        pass
    return _bare(sh, color, w)


# --------------------------------------------------------------------------- #
# Timón de Kubernetes  (marca)
# --------------------------------------------------------------------------- #
def k8s_wheel(slide, cx, cy, r, color, w=2.0, spokes=True):
    verts = [(cx + r * math.cos(-math.pi / 2 + i * 2 * math.pi / 7),
              cy + r * math.sin(-math.pi / 2 + i * 2 * math.pi / 7))
             for i in range(7)]
    poly(slide, verts, color, w)
    hub = r * 0.30
    ring(slide, cx, cy, hub, color, w)
    if spokes:
        for vx, vy in verts:
            ux, uy = (vx - cx) / r, (vy - cy) / r
            line(slide, cx + ux * hub, cy + uy * hub,
                 cx + ux * r * 0.98, cy + uy * r * 0.98, color, max(1.0, w * 0.7))


# --------------------------------------------------------------------------- #
# Set de iconos de tema
# --------------------------------------------------------------------------- #
def _w(s, w):
    return float(w) if w else max(1.2, min(3.0, s / 150000))


def i_pod(slide, x, y, s, c, w=None):
    w = _w(s, w)
    rrect(slide, x + s * .06, y + s * .20, s * .88, s * .60, c, w, rad=.28)
    ring(slide, x + s * .36, y + s * .50, s * .085, c, w)
    ring(slide, x + s * .64, y + s * .50, s * .085, c, w)


def i_wheel(slide, x, y, s, c, w=None):
    k8s_wheel(slide, x + s / 2, y + s / 2, s * .46, c, _w(s, w))


def i_doc(slide, x, y, s, c, w=None):
    w = _w(s, w)
    fold = s * .24
    poly(slide, [(x + s * .16, y + s * .04), (x + s * .70, y + s * .04),
                 (x + s * .86, y + s * .20), (x + s * .86, y + s * .96),
                 (x + s * .16, y + s * .96)], c, w)
    line(slide, x + s * .70, y + s * .04, x + s * .70, y + s * .20, c, w)
    line(slide, x + s * .70, y + s * .20, x + s * .86, y + s * .20, c, w)
    for k in (.42, .58, .74):
        line(slide, x + s * .28, y + s * k, x + s * .72, y + s * k, c, w)


def i_graph(slide, x, y, s, c, w=None):
    w = _w(s, w)
    a = (x + s * .16, y + s * .22)
    b = (x + s * .82, y + s * .30)
    d = (x + s * .5, y + s * .84)
    line(slide, *a, *b, c, w)
    line(slide, *a, *d, c, w)
    line(slide, *b, *d, c, w)
    for p in (a, b, d):
        _bare(slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(p[0] - s * .10)),
              Emu(int(p[1] - s * .10)), Emu(int(s * .20)), Emu(int(s * .20))),
              c, w, fill=D.BG)


def i_crosshair(slide, x, y, s, c, w=None):
    w = _w(s, w)
    ring(slide, x + s / 2, y + s / 2, s * .40, c, w)
    line(slide, x + s / 2, y + s * .02, x + s / 2, y + s * .24, c, w)
    line(slide, x + s / 2, y + s * .76, x + s / 2, y + s * .98, c, w)
    line(slide, x + s * .02, y + s / 2, x + s * .24, y + s / 2, c, w)
    line(slide, x + s * .76, y + s / 2, x + s * .98, y + s / 2, c, w)
    dot(slide, x + s / 2, y + s / 2, s * .07, c)


def i_cylinder(slide, x, y, s, c, w=None):
    w = _w(s, w)
    rx, ry = s * .38, s * .13
    cx = x + s / 2
    top, bot = y + s * .18, y + s * .82
    line(slide, cx - rx, top, cx - rx, bot, c, w)
    line(slide, cx + rx, top, cx + rx, bot, c, w)
    for cyv in (bot, top):                          # elipse inferior, luego superior
        _bare(slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(cx - rx)),
              Emu(int(cyv - ry)), Emu(int(2 * rx)), Emu(int(2 * ry))), c, w)


def i_stack(slide, x, y, s, c, w=None):
    w = _w(s, w)
    for k in (0, 1, 2):
        rrect(slide, x + s * .10 + k * s * .06, y + s * .66 - k * s * .28,
              s * .68, s * .20, c, w, rad=.35)


def i_cluster(slide, x, y, s, c, w=None):
    w = _w(s, w)
    rrect(slide, x + s * .04, y + s * .10, s * .92, s * .80, c, w, rad=.12)
    for gx in (.18, .44, .70):
        rrect(slide, x + s * gx, y + s * .34, s * .18, s * .32, c, w, rad=.22)


def i_key(slide, x, y, s, c, w=None):
    w = _w(s, w)
    ring(slide, x + s * .30, y + s * .34, s * .22, c, w)
    line(slide, x + s * .44, y + s * .48, x + s * .90, y + s * .90, c, w)
    line(slide, x + s * .74, y + s * .74, x + s * .88, y + s * .60, c, w)
    line(slide, x + s * .62, y + s * .62, x + s * .78, y + s * .46, c, w)


def i_shield(slide, x, y, s, c, w=None):
    w = _w(s, w)
    hw = s * .40
    poly(slide, [(x + s / 2, y + s * .04), (x + s / 2 + hw, y + s * .22),
                 (x + s / 2 + hw, y + s * .56), (x + s / 2, y + s * .96),
                 (x + s / 2 - hw, y + s * .56), (x + s / 2 - hw, y + s * .22)],
         c, w)
    line(slide, x + s * .34, y + s * .48, x + s * .46, y + s * .62, c, w)
    line(slide, x + s * .46, y + s * .62, x + s * .68, y + s * .34, c, w)


def i_cycle(slide, x, y, s, c, w=None):
    w = _w(s, w)
    r = s * .36
    cx, cy = x + s / 2, y + s / 2
    ring(slide, cx, cy, r, c, w)
    poly(slide, [(cx + r + s * .10, cy - s * .10), (cx + r - s * .12, cy - s * .02),
                 (cx + r + s * .02, cy + s * .16)], c, w, fill=c)
    poly(slide, [(cx - r - s * .10, cy + s * .10), (cx - r + s * .12, cy + s * .02),
                 (cx - r - s * .02, cy - s * .16)], c, w, fill=c)


def i_gear(slide, x, y, s, c, w=None):
    w = _w(s, w)
    cx, cy, r = x + s / 2, y + s / 2, s * .30
    ring(slide, cx, cy, r, c, w)
    dot(slide, cx, cy, s * .09, c)
    for i in range(8):
        a = i * math.pi / 4
        line(slide, cx + math.cos(a) * r, cy + math.sin(a) * r,
             cx + math.cos(a) * r * 1.42, cy + math.sin(a) * r * 1.42, c, w * 1.4)


def i_bars(slide, x, y, s, c, w=None):
    w = _w(s, w)
    line(slide, x + s * .10, y + s * .92, x + s * .92, y + s * .92, c, w)
    for gx, h in ((.22, .34), (.46, .60), (.70, .82)):
        rrect(slide, x + s * gx, y + s * .92 - s * h, s * .16, s * h, c, w,
              rad=.10)


def i_gateway(slide, x, y, s, c, w=None):
    w = _w(s, w)
    poly(slide, [(x + s * .10, y + s * .92), (x + s * .24, y + s * .30),
                 (x + s * .40, y + s * .12), (x + s * .60, y + s * .12),
                 (x + s * .76, y + s * .30), (x + s * .90, y + s * .92)],
         c, w, close=False)
    line(slide, x + s * .34, y + s * .58, x + s * .66, y + s * .58, c, w,
         arrow=True)


def i_magnifier(slide, x, y, s, c, w=None):
    w = _w(s, w)
    ring(slide, x + s * .40, y + s * .40, s * .30, c, w)
    line(slide, x + s * .62, y + s * .62, x + s * .92, y + s * .92, c, w * 1.6)


def i_clock(slide, x, y, s, c, w=None):
    w = _w(s, w)
    ring(slide, x + s / 2, y + s / 2, s * .42, c, w)
    line(slide, x + s / 2, y + s / 2, x + s / 2, y + s * .22, c, w)
    line(slide, x + s / 2, y + s / 2, x + s * .70, y + s * .58, c, w)


def i_target(slide, x, y, s, c, w=None):
    w = _w(s, w)
    ring(slide, x + s / 2, y + s / 2, s * .42, c, w)
    ring(slide, x + s / 2, y + s / 2, s * .24, c, w)
    dot(slide, x + s / 2, y + s / 2, s * .08, c)


def i_check(slide, x, y, s, c, w=None):
    w = _w(s, w)
    rrect(slide, x + s * .10, y + s * .10, s * .80, s * .80, c, w, rad=.20)
    line(slide, x + s * .30, y + s * .52, x + s * .44, y + s * .68, c, w * 1.3)
    line(slide, x + s * .44, y + s * .68, x + s * .72, y + s * .32, c, w * 1.3)


def i_bulb(slide, x, y, s, c, w=None):
    w = _w(s, w)
    ring(slide, x + s / 2, y + s * .40, s * .30, c, w)
    line(slide, x + s * .40, y + s * .74, x + s * .60, y + s * .74, c, w)
    line(slide, x + s * .42, y + s * .86, x + s * .58, y + s * .86, c, w)


def i_warn(slide, x, y, s, c, w=None):
    w = _w(s, w)
    poly(slide, [(x + s / 2, y + s * .08), (x + s * .94, y + s * .88),
                 (x + s * .06, y + s * .88)], c, w)
    line(slide, x + s / 2, y + s * .38, x + s / 2, y + s * .62, c, w * 1.3)
    dot(slide, x + s / 2, y + s * .74, s * .05, c)


def i_beaker(slide, x, y, s, c, w=None):
    w = _w(s, w)
    poly(slide, [(x + s * .34, y + s * .06), (x + s * .34, y + s * .40),
                 (x + s * .10, y + s * .92), (x + s * .90, y + s * .92),
                 (x + s * .66, y + s * .40), (x + s * .66, y + s * .06)],
         c, w, close=False)
    line(slide, x + s * .28, y + s * .06, x + s * .72, y + s * .06, c, w)
    line(slide, x + s * .24, y + s * .66, x + s * .76, y + s * .66, c, w)


def i_bug(slide, x, y, s, c, w=None):
    w = _w(s, w)
    ring(slide, x + s / 2, y + s * .30, s * .16, c, w)
    _bare(slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(x + s * .28)),
          Emu(int(y + s * .40)), Emu(int(s * .44)), Emu(int(s * .50))), c, w)
    for sy in (.5, .66, .82):
        line(slide, x + s * .30, y + s * sy, x + s * .06, y + s * (sy - .06),
             c, w)
        line(slide, x + s * .70, y + s * sy, x + s * .94, y + s * (sy - .06),
             c, w)


def i_list(slide, x, y, s, c, w=None):
    w = _w(s, w)
    for k in (.24, .48, .72):
        dot(slide, x + s * .16, y + s * k, s * .06, c)
        line(slide, x + s * .32, y + s * k, x + s * .90, y + s * k, c, w)


def i_bolt(slide, x, y, s, c, w=None):
    w = _w(s, w)
    poly(slide, [(x + s * .54, y + s * .04), (x + s * .20, y + s * .56),
                 (x + s * .46, y + s * .56), (x + s * .40, y + s * .96),
                 (x + s * .80, y + s * .40), (x + s * .52, y + s * .40)],
         c, w)


ICONS = {
    "pod": i_pod, "wheel": i_wheel, "doc": i_doc, "graph": i_graph,
    "crosshair": i_crosshair, "cylinder": i_cylinder, "stack": i_stack,
    "cluster": i_cluster, "key": i_key, "shield": i_shield, "cycle": i_cycle,
    "gear": i_gear, "bars": i_bars, "gateway": i_gateway,
    "magnifier": i_magnifier, "clock": i_clock, "target": i_target,
    "check": i_check, "bulb": i_bulb, "warn": i_warn, "beaker": i_beaker,
    "bug": i_bug, "list": i_list, "bolt": i_bolt,
}


# --------------------------------------------------------------------------- #
# Inferencia por palabras clave
# --------------------------------------------------------------------------- #
_KEYMAP = [
    # --- iconos oficiales de Kubernetes (PNG) ---
    (("networkpolicy", "network policy", "firewall", "aislar", "default deny", "la red está rota"), "netpol"),
    (("rbac", "forbidden", "permiso", "privilegio", "auth can"), "role"),
    (("clusterip", "nodeport", "loadbalancer", "externalname", "headless"), "svc"),
    (("statefulset", "stateful"), "sts"),
    (("daemonset",), "ds"),
    (("replicaset", "replicationcontroller", "self-healing"), "rs"),
    (("workload", "controlador", "deployment", "réplica", "despliegue", "rolling update", "rollback", "rollout", "lifecycle", "revisi"), "deploy"),
    (("integrador", "arquitectura a construir", "todo el curso", "frontend", "backend"), "deploy"),
    (("etcd", "snapshot", "backup", "restore"), "etcd"),
    (("secret",), "secret"),
    (("configmap", "config", "externaliza"), "cm"),
    (("storage", "almacenamiento", "volumen", "volume", "persistentvolume", " pv", "pvc", "emptydir", "statefulsets"), "pv"),
    (("scheduling", "scheduler", "colocación", "nodeselector", "affinity", "taint", "priorityclass"), "sched"),
    (("ingress", "gateway", "tls", "https", "puerta de entrada"), "ing"),
    (("service", "endpointslice", "coredns", " dns", "networking", "tráfico", "ruta del tráfico", "del cliente al contenedor"), "svc"),
    (("mental model", "relación", "cadena", "ruta de", "seis mental", "cada capa"), "svc"),
    (("kubeadm", "control plane", "cluster", "static pod", "nodo y control", "no responde"), "controlplane"),
    (("recurso", "límit", "limit", "request", "qos", "métrica", "metrics", "escalado", "hpa"), "limits"),
    (("namespace",), "ns"),
    (("pod",), "pod"),
    (("kubernetes",), "wheel"),
    # --- iconos vectoriales (conceptos sin recurso K8s) ---
    (("troubleshoot", "diagnóst", "diagnostic", "el método", "método general", "método", "seis pasos"), "magnifier"),
    (("upgrade",), "cycle"),
    (("objetivo",), "target"),
    (("agenda", "minutos"), "clock"),
    (("checklist", "comprueba"), "check"),
    (("tip", "atajo", "devuelven minutos"), "bulb"),
    (("aviso", "actualización técnica", "corrección", "matiz crítico", "fallo estrella", "guion lo cambia", "no reequilibra"), "warn"),
    (("demo", "demostración"), "bolt"),
    (("sintaxis", "manifiesto", "yaml", "fragmentos"), "doc"),
    (("estrategia de examen", "orden de construcción", "los tres comandos", "los cinco comandos", "procedimiento", "secuencia"), "list"),
]


def name_for(*texts, default="wheel"):
    blob = " ".join(t for t in texts if t).lower()
    for keys, icon in _KEYMAP:
        if any(k in blob for k in keys):
            return icon
    return default


def draw(slide, icon, x, y, s, color, w=None):
    """Dibuja `icon`: PNG oficial de Kubernetes si existe, si no, línea vectorial."""
    p = _png_path(icon)
    if p:
        slide.shapes.add_picture(p, Emu(int(x)), Emu(int(y)),
                                 Emu(int(s)), Emu(int(s)))
        return
    ICONS.get(icon, i_wheel)(slide, x, y, s, color, w)


def draw_named(slide, x, y, s, color, *texts, default="wheel", w=None):
    draw(slide, name_for(*texts, default=default), x, y, s, color, w)
