"""
Renderizadores de arquetipos de diapositiva y ensamblado del .pptx.

El contenido de cada clase se describe como una lista de dicts (ver
`contenido/clase01.py`). Cada dict tiene una clave `t` con el arquetipo.
"""

import copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import design as D
from design import (
    EMU_W, EMU_H, MARGIN, CONTENT_W, COL_GAP, Y_BODY, Y_FOOTER,
    INK, BODY, MUTED, FAINT, PANEL, PANEL_SOFT, HAIRLINE, BG, BG_ALT,
    K8S, K8S_TEXT, K8S_DEEP, ACCENTS, CODE_BG,
    F_HEAD, F_LIGHT, F_BODY, F_MONO,
    SZ_H1, SZ_H_SECTION, SZ_LEAD, SZ_BODY, SZ_SMALL, SZ_EYEBROW,
)

BODY_H = Y_FOOTER - 120000 - Y_BODY


def _accent(name):
    return ACCENTS.get(name or "blue", ACCENTS["blue"])


# --------------------------------------------------------------------------- #
# Arquetipos
# --------------------------------------------------------------------------- #
def render_cover(slide, s, ctx):
    D.set_bg(slide, BG)
    # motivos geométricos
    D.oval(slide, EMU_W - 3050000, -1350000, 4500000, 4500000,
           fill=D.BG_ALT)
    D.oval(slide, EMU_W - 1850000, EMU_H - 1550000, 2900000, 2900000,
           fill=D.PANEL_SOFT)
    D.hex_motif(slide, cx=EMU_W - 1350000, cy=1550000, r=1150000,
                color=D.HAIRLINE, line_w=2.25)

    x = MARGIN + 40000
    D.eyebrow(slide, ctx["curso"], x=x, w=8600000)
    tb = D.textbox(slide, x, 1620000, 8600000, 520000)
    D.para(tb.text_frame, s["kicker"], size=15, color=K8S_TEXT, font=F_HEAD,
           bold=True, spacing=2.0, first=True, space_after=0, line=1.0)
    tb = D.textbox(slide, x, 2040000, 8900000, 1500000)
    D.para(tb.text_frame, s["title"], size=44, color=INK, font=F_LIGHT,
           first=True, space_after=0, line=1.06)
    D.accent_rule(slide, x=x, y=3560000, w=620000, color=K8S)
    tb = D.textbox(slide, x, 3760000, 8600000, 520000)
    D.para(tb.text_frame, s["subtitle"], size=15, color=BODY, font=F_BODY,
           italic=True, first=True, space_after=0, line=1.2)

    # tira de estadísticas
    stats = s["stats"]
    n = len(stats)
    gap = 150000
    total = CONTENT_W + 80000
    cw = (total - gap * (n - 1)) / n
    y = 4650000
    for i, (val, lab) in enumerate(stats):
        cx = x + i * (cw + gap)
        D.rect(slide, cx, y, cw, 900000, fill=PANEL, line=HAIRLINE,
               line_w=1.0, rounded=True, radius=0.10)
        D.rect(slide, cx, y, cw, 46000, fill=K8S, line=None)
        c = D.textbox(slide, cx + 150000, y + 120000, cw - 300000, 660000,
                      anchor=MSO_ANCHOR.MIDDLE)
        D.para(c.text_frame, str(val), size=22, color=INK, font=F_LIGHT,
               first=True, space_after=2, line=1.0)
        D.para(c.text_frame, lab.upper(), size=SZ_SMALL - 1, color=MUTED,
               font=F_HEAD, bold=True, spacing=1.6, space_after=0, line=1.0)


def render_section(slide, s, ctx):
    ac, ac_t, ac_d = _accent(s.get("accent"))
    D.set_bg(slide, BG_ALT)
    D.oval(slide, -1500000, EMU_H - 2600000, 4200000, 4200000, fill=BG)
    D.hex_motif(slide, cx=EMU_W - 1250000, cy=EMU_H - 1150000, r=1500000,
                color=D.HAIRLINE, line_w=2.5)

    # número fantasma
    tb = D.textbox(slide, MARGIN - 30000, 1150000, 4000000, 2600000)
    D.para(tb.text_frame, f"{int(s['index']):02d}", size=150, color=FAINT,
           font=F_LIGHT, first=True, space_after=0, line=0.9)

    tb = D.textbox(slide, MARGIN + 20000, 3050000, 9200000, 320000,
                   anchor=MSO_ANCHOR.MIDDLE)
    D.para(tb.text_frame, s.get("eyebrow", f"BLOQUE {int(s['index'])}").upper(),
           size=SZ_EYEBROW, color=ac_t, font=F_HEAD, bold=True, spacing=2.6,
           first=True, space_after=0, line=1.0)
    tb = D.textbox(slide, MARGIN + 20000, 3350000, 10200000, 1000000)
    D.para(tb.text_frame, s["title"], size=SZ_H_SECTION, color=INK,
           font=F_LIGHT, first=True, space_after=0, line=1.02)
    D.rect(slide, MARGIN + 24000, 4380000, 620000, 46000, fill=ac, line=None)
    tb = D.textbox(slide, MARGIN + 20000, 4560000, 9600000, 620000)
    D.para(tb.text_frame, s["subtitle"], size=SZ_LEAD, color=BODY, font=F_BODY,
           italic=True, first=True, space_after=0, line=1.2)
    D.footer(slide, ctx["clase_n"], ctx["page"])


def render_list(slide, s, ctx):
    ac, ac_t, ac_d = _accent(s.get("accent"))
    D.scaffold(slide, clase_n=ctx["clase_n"], page=ctx["page"],
               eyebrow_text=s["eyebrow"], heading_text=s["heading"],
               accent_text=ac_t, rule_color=ac)
    items = s["items"]
    cols = s.get("columns", 1)
    check = s.get("style") == "check"
    per = (len(items) + cols - 1) // cols
    col_w = (CONTENT_W - (COL_GAP if cols == 2 else 0)) / cols
    d = 300000
    for c in range(cols):
        group = items[c * per:(c + 1) * per]
        x = MARGIN + c * (col_w + COL_GAP)
        row_h = min(560000, BODY_H / max(per, 1))
        y0 = Y_BODY + (BODY_H - row_h * len(group)) / 2
        for i, it in enumerate(group):
            n = c * per + i + 1
            y = y0 + i * row_h
            if check:
                bx = D.rect(slide, x, y + (row_h - d) / 2, d, d, fill=None,
                            line=ac, line_w=1.5, rounded=True, radius=0.28)
                tf = bx.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                D.para(tf, "✓", size=13, color=ac_t, font=F_BODY, bold=True,
                       align=PP_ALIGN.CENTER, first=True, space_after=0, line=1.0)
            else:
                D.num_badge(slide, x, y + (row_h - d) / 2, d, n, fill=ac)
            tb = D.textbox(slide, x + d + 200000, y, col_w - d - 200000, row_h,
                           anchor=MSO_ANCHOR.MIDDLE)
            D.para(tb.text_frame, it, size=SZ_BODY, color=BODY, font=F_BODY,
                   first=True, space_after=0, line=1.16)


def render_agenda(slide, s, ctx):
    D.scaffold(slide, clase_n=ctx["clase_n"], page=ctx["page"],
               eyebrow_text=s["eyebrow"], heading_text=s["heading"])
    rows = s["rows"]
    row_h = min(470000, BODY_H / len(rows))
    y0 = Y_BODY + (BODY_H - row_h * len(rows)) / 2
    tw = 1500000
    for i, (t, txt) in enumerate(rows):
        y = y0 + i * row_h
        lab = any(k in txt for k in ("LAB", "Lab"))
        if lab:
            D.rect(slide, MARGIN - 60000, y + 40000, CONTENT_W + 120000,
                   row_h - 80000, fill=PANEL_SOFT, line=None, rounded=True,
                   radius=0.16)
        tb = D.textbox(slide, MARGIN, y, tw, row_h, anchor=MSO_ANCHOR.MIDDLE)
        D.para(tb.text_frame, t, size=SZ_SMALL, color=K8S_TEXT if lab else MUTED,
               font=F_MONO, bold=True, first=True, space_after=0, line=1.0,
               spacing=0.4)
        D.rect(slide, MARGIN + tw, y + row_h / 2 - 60000, 8000, 120000,
               fill=K8S if lab else HAIRLINE, line=None)
        tb = D.textbox(slide, MARGIN + tw + 200000, y,
                       CONTENT_W - tw - 200000, row_h, anchor=MSO_ANCHOR.MIDDLE)
        D.para(tb.text_frame, txt, size=SZ_BODY,
               color=INK if lab else BODY, font=F_BODY, bold=lab,
               first=True, space_after=0, line=1.1)


def render_concept(slide, s, ctx):
    ac, ac_t, ac_d = _accent(s.get("accent"))
    D.scaffold(slide, clase_n=ctx["clase_n"], page=ctx["page"],
               eyebrow_text=s["eyebrow"], heading_text=s["heading"],
               accent_text=ac_t, rule_color=ac)
    left_w = CONTENT_W * 0.52
    right_x = MARGIN + left_w + COL_GAP
    right_w = CONTENT_W - left_w - COL_GAP
    note = s.get("note")
    avail_h = BODY_H - (600000 if note else 0)

    pts = s["points"]
    ph = min(avail_h / len(pts), 940000)
    block_h = ph * len(pts)
    y0 = Y_BODY + (avail_h - block_h) / 2
    for i, p in enumerate(pts):
        y = y0 + i * ph
        D.oval(slide, MARGIN + 16000, y + ph / 2 - 52000, 104000, 104000,
               fill=ac)
        tb = D.textbox(slide, MARGIN + 260000, y, left_w - 260000, ph,
                       anchor=MSO_ANCHOR.MIDDLE)
        D.para(tb.text_frame, p, size=SZ_BODY, color=BODY, font=F_BODY,
               first=True, space_after=0, line=1.18)

    side = s["side"]
    if side["kind"] == "code":
        D.code_panel(slide, right_x, y0, right_w, side["lines"],
                     accent=ac, title=side.get("title"), min_h=block_h,
                     size=SZ_SMALL)
    else:
        line_h = 300000
        nat = 220000 + 300000 + len(side["lines"]) * line_h + 240000
        panel_h = max(block_h, min(avail_h, nat))
        D.rect(slide, right_x, y0, right_w, panel_h, fill=PANEL,
               line=HAIRLINE, line_w=1.0, rounded=True, radius=0.05)
        D.rect(slide, right_x + 24000, y0, right_w - 48000, 46000, fill=ac,
               line=None)
        tb = D.textbox(slide, right_x + 240000, y0 + 240000,
                       right_w - 480000, panel_h - 480000,
                       anchor=MSO_ANCHOR.MIDDLE)
        tf = tb.text_frame
        D.para(tf, side["title"].upper(), size=SZ_SMALL - 1, color=ac_t,
               font=F_HEAD, bold=True, spacing=1.8, first=True, space_after=14,
               line=1.0)
        for ln in side["lines"]:
            D.para(tf, ln, size=SZ_SMALL, color=BODY, font=F_BODY,
                   space_after=11, line=1.2, bullet=True)

    if note:
        nac = _accent(note.get("accent") or s.get("accent"))
        D.note_band(slide, note["text"], y=Y_BODY + avail_h + 150000,
                    accent=nac[0], label=note.get("label"), label_color=nac[1])


def render_columns(slide, s, ctx):
    ac, ac_t, ac_d = _accent(s.get("accent"))
    D.scaffold(slide, clase_n=ctx["clase_n"], page=ctx["page"],
               eyebrow_text=s["eyebrow"], heading_text=s["heading"],
               accent_text=ac_t, rule_color=ac)
    cols = s["cols"]
    note = s.get("note")
    avail_h = BODY_H - (600000 if note else 0)
    n = len(cols)
    cw = (CONTENT_W - COL_GAP * (n - 1)) / n
    pad = 260000
    cpl = max(24, int((cw - 2 * pad) / 78000))

    def wrapped(lines):
        return sum(max(1, -(-len(t) // cpl)) for t in lines)

    has_tag = any(c.get("tag") for c in cols)
    has_sub = any(c.get("subtitle") for c in cols)
    hh = 320000 + (240000 if has_tag else 0) + (250000 if has_sub else 0)
    body_rows = max(wrapped(c["lines"]) for c in cols)
    panel_h = min(avail_h, 260000 + hh + 200000 + body_rows * 288000 + 300000)
    top = Y_BODY + (avail_h - panel_h) / 2
    for i, col in enumerate(cols):
        x = MARGIN + i * (cw + COL_GAP)
        cac = _accent(col.get("accent") or s.get("accent"))
        D.rect(slide, x, top, cw, panel_h, fill=PANEL, line=HAIRLINE,
               line_w=1.0, rounded=True, radius=0.05)
        D.rect(slide, x + 24000, top, cw - 48000, 46000, fill=cac[0], line=None)
        tb = D.textbox(slide, x + pad, top + 260000, cw - 2 * pad, hh)
        tf = tb.text_frame
        if col.get("tag"):
            D.para(tf, col["tag"].upper(), size=SZ_SMALL - 1, color=cac[1],
                   font=F_HEAD, bold=True, spacing=1.8, first=True,
                   space_after=7, line=1.0)
        D.para(tf, col["title"], size=SZ_LEAD + 1, color=INK, font=F_HEAD,
               bold=True, first=not col.get("tag"), space_after=4, line=1.05)
        if col.get("subtitle"):
            D.para(tf, col["subtitle"], size=SZ_SMALL, color=MUTED, font=F_BODY,
                   italic=True, space_after=0, line=1.12)
        tb = D.textbox(slide, x + pad, top + 260000 + hh,
                       cw - 2 * pad, panel_h - 260000 - hh - 180000)
        tf = tb.text_frame
        for ln in col["lines"]:
            st = ln.strip()
            mono = (st.startswith(("$ ", "k ", "kubectl ", "- ", "-", "#"))
                    or ln.startswith("  ") or ": {" in ln or st.endswith(":"))
            D.para(tf, ln, size=SZ_SMALL - (1 if mono else 0), color=BODY,
                   font=F_MONO if mono else F_BODY,
                   space_after=6 if mono else 8, line=1.3 if mono else 1.16,
                   bullet=not mono)
    if note:
        nac = _accent(note.get("accent") or s.get("accent"))
        D.note_band(slide, note["text"], y=Y_BODY + avail_h + 150000,
                    accent=nac[0], label=note.get("label"), label_color=nac[1])


def render_codeblocks(slide, s, ctx):
    ac, ac_t, ac_d = _accent(s.get("accent"))
    D.scaffold(slide, clase_n=ctx["clase_n"], page=ctx["page"],
               eyebrow_text=s["eyebrow"], heading_text=s["heading"],
               accent_text=ac_t, rule_color=ac)
    blocks = s["blocks"]
    note = s.get("note")
    intro = s.get("intro")
    y = Y_BODY
    avail = BODY_H
    if intro:
        tb = D.textbox(slide, MARGIN, y, CONTENT_W, 520000)
        D.para(tb.text_frame, intro, size=SZ_BODY, color=BODY, font=F_BODY,
               italic=True, first=True, space_after=0, line=1.16)
        y += 560000
        avail -= 560000
    if note:
        avail -= 560000
    gap = 200000
    pad_v = 300000

    def layout(sz):
        lh = int(224000 * sz / SZ_SMALL)
        h = [pad_v + (300000 if b.get("label") else 0) + len(b["lines"]) * lh
             for b in blocks]
        return h, sum(h) + gap * (len(blocks) - 1)

    # baja el tamaño de fuente del código hasta que todos los bloques quepan
    code_sz = SZ_SMALL
    heights, total = layout(code_sz)
    while total > avail and code_sz > 8.5:
        code_sz -= 0.5
        heights, total = layout(code_sz)

    if total > avail:                     # aún no cabe: escala como último recurso
        heights = [h * avail / total for h in heights]
        y0 = y
    else:
        y0 = y + min((avail - total) / 2, 380000)
    offs = [y0 + sum(heights[:i]) + gap * i for i in range(len(blocks))]
    for i, b in enumerate(blocks):
        by = offs[i]
        bh = heights[i]
        lines = b["lines"]
        D.rect(slide, MARGIN, by, CONTENT_W, bh, fill=CODE_BG, line=HAIRLINE,
               line_w=1.0, rounded=True, radius=0.05)
        D.rect(slide, MARGIN, by, 46000, bh, fill=ac, line=None)
        if b.get("label"):
            tb = D.textbox(slide, MARGIN + 230000, by + 120000,
                           CONTENT_W - 460000, 260000)
            D.para(tb.text_frame, b["label"].upper(), size=SZ_SMALL - 1,
                   color=MUTED, font=F_HEAD, bold=True, spacing=1.8, first=True,
                   space_after=0, line=1.0)
        tb = D.textbox(slide, MARGIN + 270000, by + (400000 if b.get("label") else 150000),
                       CONTENT_W - 540000, bh - (460000 if b.get("label") else 260000))
        tf = tb.text_frame
        for j, ln in enumerate(lines):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.line_spacing = 1.34
            p.space_after = Pt(0)
            D._set_bullet(p, False)
            D._emit_code_line(p, ln, code_sz)
    if note:
        nac = _accent(note.get("accent") or s.get("accent"))
        D.note_band(slide, note["text"], y=Y_FOOTER - 120000 - 470000,
                    accent=nac[0], label=note.get("label"),
                    label_color=nac[1])


def render_table(slide, s, ctx):
    ac, ac_t, ac_d = _accent(s.get("accent"))
    D.scaffold(slide, clase_n=ctx["clase_n"], page=ctx["page"],
               eyebrow_text=s["eyebrow"], heading_text=s["heading"],
               accent_text=ac_t, rule_color=ac)
    headers = s["headers"]
    rows = s["rows"]
    note = s.get("note")
    avail_h = BODY_H - (520000 if note else 0)
    ncol = len(headers)
    weights = s.get("weights", [1] * ncol)
    tot = sum(weights)
    widths = [CONTENT_W * w / tot for w in weights]
    head_h = 420000
    row_h = min(560000, (avail_h - head_h) / len(rows))
    x0, y0 = MARGIN, Y_BODY
    # cabecera
    D.rect(slide, x0, y0, CONTENT_W, head_h, fill=PANEL, line=None,
           rounded=False)
    D.rect(slide, x0, y0 + head_h - 8000, CONTENT_W, 10000, fill=ac, line=None)
    cx = x0
    for j, htxt in enumerate(headers):
        tb = D.textbox(slide, cx + 160000, y0, widths[j] - 220000, head_h,
                       anchor=MSO_ANCHOR.MIDDLE)
        D.para(tb.text_frame, htxt.upper(), size=SZ_SMALL - 1, color=ac_t,
               font=F_HEAD, bold=True, spacing=1.6, first=True, space_after=0,
               line=1.0)
        cx += widths[j]
    # filas
    for i, row in enumerate(rows):
        ry = y0 + head_h + i * row_h
        if i % 2 == 1:
            D.rect(slide, x0, ry, CONTENT_W, row_h, fill=PANEL_SOFT, line=None)
        D.rect(slide, x0, ry + row_h, CONTENT_W, 8000, fill=HAIRLINE, line=None)
        cx = x0
        for j, cell in enumerate(row):
            first_col = j == 0
            mono = j == ncol - 1 and s.get("mono_last", True)
            tb = D.textbox(slide, cx + 160000, ry, widths[j] - 220000, row_h,
                           anchor=MSO_ANCHOR.MIDDLE)
            D.para(tb.text_frame, cell, size=SZ_SMALL,
                   color=INK if first_col else BODY,
                   font=F_MONO if mono else F_BODY,
                   bold=first_col, first=True, space_after=0, line=1.1)
            cx += widths[j]
    if note:
        nac = _accent(note.get("accent") or s.get("accent"))
        D.note_band(slide, note["text"], y=Y_BODY + avail_h + 120000,
                    accent=nac[0], label=note.get("label"), label_color=nac[1])


def render_cards(slide, s, ctx):
    ac, ac_t, ac_d = _accent(s.get("accent"))
    D.scaffold(slide, clase_n=ctx["clase_n"], page=ctx["page"],
               eyebrow_text=s["eyebrow"], heading_text=s["heading"],
               accent_text=ac_t, rule_color=ac)
    cards = s["cards"]
    note = s.get("note")
    ncols = s.get("cols", 2)
    avail_h = BODY_H - (600000 if note else 0)
    nrows = (len(cards) + ncols - 1) // ncols
    gx, gy = COL_GAP, 240000
    cw = (CONTENT_W - gx * (ncols - 1)) / ncols
    pad = 230000
    # altura natural: barra + título + cuerpo estimado
    cpl = max(20, int((cw - 2 * pad) / 82000))       # caracteres por línea aprox.
    est_lines = max((-(-len(c["body"]) // cpl)) for c in cards)
    nat_ch = pad + 400000 + est_lines * 236000 + pad
    ch = min((avail_h - gy * (nrows - 1)) / nrows, nat_ch, 2050000)
    grid_h = ch * nrows + gy * (nrows - 1)
    y_top = Y_BODY + (avail_h - grid_h) / 2
    for idx, card in enumerate(cards):
        r, c = divmod(idx, ncols)
        x = MARGIN + c * (cw + gx)
        y = y_top + r * (ch + gy)
        D.rect(slide, x, y, cw, ch, fill=PANEL, line=HAIRLINE, line_w=1.0,
               rounded=True, radius=0.06)
        D.rect(slide, x + 24000, y, cw - 48000, 46000, fill=ac, line=None)
        D.num_badge(slide, x + pad, y + pad + 20000, 300000, idx + 1, fill=ac)
        tb = D.textbox(slide, x + pad + 410000, y + pad, cw - pad - 460000,
                       340000, anchor=MSO_ANCHOR.MIDDLE)
        D.para(tb.text_frame, card["title"], size=SZ_LEAD, color=INK,
               font=F_HEAD, bold=True, first=True, space_after=0, line=1.05)
        tb = D.textbox(slide, x + pad, y + pad + 420000, cw - 2 * pad,
                       ch - pad - 500000)
        D.para(tb.text_frame, card["body"], size=SZ_SMALL, color=BODY,
               font=F_BODY, first=True, space_after=0, line=1.18)
    if note:
        nac = _accent(note.get("accent") or s.get("accent"))
        D.note_band(slide, note["text"], y=Y_BODY + avail_h + 150000,
                    accent=nac[0], label=note.get("label"), label_color=nac[1])


def render_process(slide, s, ctx):
    ac, ac_t, ac_d = _accent(s.get("accent"))
    D.scaffold(slide, clase_n=ctx["clase_n"], page=ctx["page"],
               eyebrow_text=s["eyebrow"], heading_text=s["heading"],
               accent_text=ac_t, rule_color=ac)
    steps = s["steps"]
    note = s.get("note")
    avail_h = BODY_H - (520000 if note else 0)
    n = len(steps)
    gap = 170000
    sh = (avail_h - gap * (n - 1)) / n
    for i, st in enumerate(steps):
        y = Y_BODY + i * (sh + gap)
        D.rect(slide, MARGIN, y, CONTENT_W, sh, fill=PANEL, line=HAIRLINE,
               line_w=1.0, rounded=True, radius=0.08)
        D.num_badge(slide, MARGIN + 220000, y + sh / 2 - 165000, 330000,
                    i + 1, fill=ac)
        tb = D.textbox(slide, MARGIN + 750000, y, 2500000, sh,
                       anchor=MSO_ANCHOR.MIDDLE)
        D.para(tb.text_frame, st["title"], size=SZ_LEAD, color=INK, font=F_HEAD,
               bold=True, first=True, space_after=0, line=1.05)
        D.rect(slide, MARGIN + 3350000, y + sh * 0.22, 8000, sh * 0.56,
               fill=HAIRLINE, line=None)
        tb = D.textbox(slide, MARGIN + 3550000, y, CONTENT_W - 3550000 - 200000,
                       sh, anchor=MSO_ANCHOR.MIDDLE)
        D.para(tb.text_frame, st["body"], size=SZ_SMALL, color=BODY,
               font=F_BODY, first=True, space_after=0, line=1.16)
        if i < n - 1:
            ar = D.textbox(slide, MARGIN + 385000 - 90000, y + sh + 8000,
                           180000, gap - 16000, anchor=MSO_ANCHOR.MIDDLE)
            D.para(ar.text_frame, "▼", size=9, color=ac_t, font=F_BODY,
                   align=PP_ALIGN.CENTER, first=True, space_after=0, line=1.0)
    if note:
        nac = _accent(note.get("accent") or s.get("accent"))
        D.note_band(slide, note["text"], y=Y_BODY + avail_h + 120000,
                    accent=nac[0], label=note.get("label"), label_color=nac[1])


def render_chain(slide, s, ctx):
    ac, ac_t, ac_d = _accent(s.get("accent"))
    D.scaffold(slide, clase_n=ctx["clase_n"], page=ctx["page"],
               eyebrow_text=s["eyebrow"], heading_text=s["heading"],
               accent_text=ac_t, rule_color=ac)
    nodes = s["nodes"]
    detail = s["detail"]
    n = len(nodes)
    gap = 150000
    per_row = n if n <= 5 else -(-n // 2)      # 2 filas si hay más de 5
    rows = [nodes[i:i + per_row] for i in range(0, n, per_row)]
    node_h = 720000 if len(rows) == 1 else 660000
    row_gap = 150000
    y = Y_BODY + (320000 if len(rows) == 1 else 200000)
    for r, row in enumerate(rows):
        ry = y + r * (node_h + row_gap)
        nw = (CONTENT_W - gap * (per_row - 1)) / per_row
        for i, node in enumerate(row):
            x = MARGIN + i * (nw + gap)
            D.rect(slide, x, ry, nw, node_h, fill=PANEL, line=ac, line_w=1.25,
                   rounded=True, radius=0.12)
            tb = D.textbox(slide, x + 70000, ry, nw - 140000, node_h,
                           anchor=MSO_ANCHOR.MIDDLE)
            D.para(tb.text_frame, node, size=SZ_SMALL, color=INK, font=F_HEAD,
                   bold=True, align=PP_ALIGN.CENTER, first=True, space_after=0,
                   line=1.12)
            last_in_deck = (r * per_row + i) == n - 1
            if i < len(row) - 1:
                ar = D.textbox(slide, x + nw - 40000, ry, gap + 80000, node_h,
                               anchor=MSO_ANCHOR.MIDDLE)
                D.para(ar.text_frame, "→", size=15, color=ac_t, font=F_BODY,
                       align=PP_ALIGN.CENTER, first=True, space_after=0, line=1.0)
            elif not last_in_deck:               # salto de fila
                ar = D.textbox(slide, x + nw / 2 - 100000, ry + node_h - 30000,
                               200000, row_gap + 60000, anchor=MSO_ANCHOR.MIDDLE)
                D.para(ar.text_frame, "↓", size=14, color=ac_t, font=F_BODY,
                       align=PP_ALIGN.CENTER, first=True, space_after=0, line=1.0)

    dy = y + len(rows) * node_h + (len(rows) - 1) * row_gap + 320000
    max_dy = Y_FOOTER - 160000
    line_h = 250000
    detail_h = min(max_dy - dy,
                   260000 + 320000 + len(detail["lines"]) * line_h + 240000)
    D.rect(slide, MARGIN, dy, CONTENT_W, detail_h, fill=PANEL_SOFT,
           line=HAIRLINE, line_w=1.0, rounded=True, radius=0.06)
    D.rect(slide, MARGIN + 30000, dy + 24000, 40000, detail_h - 48000,
           fill=ac, line=None)
    tb = D.textbox(slide, MARGIN + 260000, dy + 240000, CONTENT_W - 520000,
                   detail_h - 400000)
    tf = tb.text_frame
    D.para(tf, detail["label"].upper(), size=SZ_SMALL - 1, color=ac_t,
           font=F_HEAD, bold=True, spacing=1.8, first=True, space_after=12,
           line=1.0)
    for ln in detail["lines"]:
        mono = ln.strip().startswith(("k ", "kubectl ", "$ ", "0/"))
        D.para(tf, ln, size=SZ_SMALL, color=BODY,
               font=F_MONO if mono else F_BODY, space_after=7, line=1.2,
               bullet=not mono)


def render_lab(slide, s, ctx):
    ac, ac_t, ac_d = _accent(s.get("accent", "blue"))
    D.set_bg(slide, BG_ALT)
    D.oval(slide, EMU_W - 2200000, -1500000, 3400000, 3400000, fill=BG)
    x = MARGIN + 20000
    chip_w = max(2500000, int(len(s["level"]) * 112000) + 460000)
    D.chip(slide, x, 470000, chip_w, 360000, s["level"], fill=None,
           line=ac, text_color=ac_t, size=SZ_SMALL - 1, spacing=1.4)
    tb = D.textbox(slide, x, 990000, CONTENT_W - 40000, 760000)
    D.para(tb.text_frame, s["title"], size=31, color=INK, font=F_LIGHT,
           first=True, space_after=0, line=1.06)
    left_w = CONTENT_W - 3500000
    tb = D.textbox(slide, x, 1800000, left_w, 640000)
    D.para(tb.text_frame, s["desc"], size=SZ_BODY, color=BODY, font=F_BODY,
           italic=True, first=True, space_after=0, line=1.2)

    steps = s["steps"]
    y0 = 2560000
    row_h = min(560000, (Y_FOOTER - 360000 - y0) / len(steps))
    for i, st in enumerate(steps):
        y = y0 + i * row_h
        D.num_badge(slide, x, y + (row_h - 300000) / 2, 300000, i + 1, fill=ac)
        tb = D.textbox(slide, x + 470000, y, left_w - 470000, row_h,
                       anchor=MSO_ANCHOR.MIDDLE)
        D.para(tb.text_frame, st, size=SZ_SMALL, color=BODY, font=F_BODY,
               first=True, space_after=0, line=1.16)

    # panel de metadatos
    mx = MARGIN + left_w + COL_GAP
    mw = CONTENT_W - left_w - COL_GAP
    my = 2560000
    mh = Y_FOOTER - 360000 - my
    D.rect(slide, mx, my, mw, mh, fill=PANEL, line=HAIRLINE, line_w=1.0,
           rounded=True, radius=0.06)
    D.rect(slide, mx, my, mw, 46000, fill=ac, line=None)
    meta = [("Archivo", s["file"]), ("Duración", s["duration"]),
            ("Validación", s["validation"])]
    ih = (mh - 240000) / len(meta)
    for i, (k, v) in enumerate(meta):
        iy = my + 180000 + i * ih
        tb = D.textbox(slide, mx + 200000, iy, mw - 400000, ih - 60000)
        tf = tb.text_frame
        D.para(tf, k.upper(), size=SZ_SMALL - 2, color=MUTED, font=F_HEAD,
               bold=True, spacing=1.8, first=True, space_after=5, line=1.0)
        D.para(tf, v, size=SZ_SMALL, color=INK,
               font=F_MONO if k != "Duración" else F_HEAD, bold=k == "Duración",
               space_after=0, line=1.12)
        if i < len(meta) - 1:
            D.rect(slide, mx + 200000, iy + ih - 40000, mw - 400000, 8000,
                   fill=HAIRLINE, line=None)
    D.footer(slide, ctx["clase_n"], ctx["page"])


def render_closing(slide, s, ctx):
    D.set_bg(slide, BG_ALT)
    D.oval(slide, EMU_W - 2600000, -1400000, 4000000, 4000000, fill=BG)
    D.hex_motif(slide, cx=1250000, cy=EMU_H - 1100000, r=1450000,
                color=D.HAIRLINE, line_w=2.5)
    x = MARGIN + 20000
    D.eyebrow(slide, s["eyebrow"], x=x, w=9000000)
    tb = D.textbox(slide, x, 1500000, 10000000, 900000)
    D.para(tb.text_frame, s["title"], size=46, color=INK, font=F_LIGHT,
           first=True, space_after=0, line=1.04)
    D.accent_rule(slide, x=x + 4000, y=2560000, w=620000, color=K8S)
    tb = D.textbox(slide, x, 2820000, CONTENT_W - 1200000, 2600000)
    tf = tb.text_frame
    for i, pgraph in enumerate(s["paragraphs"]):
        D.para(tf, pgraph, size=SZ_LEAD, color=BODY, font=F_BODY,
               first=i == 0, space_after=14, line=1.32)
    D.footer(slide, ctx["clase_n"], ctx["page"])


RENDERERS = {
    "cover": render_cover, "section": render_section, "list": render_list,
    "agenda": render_agenda, "concept": render_concept, "columns": render_columns,
    "codeblocks": render_codeblocks, "table": render_table, "cards": render_cards,
    "process": render_process, "chain": render_chain, "lab": render_lab,
    "closing": render_closing,
}


# --------------------------------------------------------------------------- #
# Ensamblado
# --------------------------------------------------------------------------- #
def build(deck, out_path):
    prs = Presentation()
    prs.slide_width = Emu(EMU_W)
    prs.slide_height = Emu(EMU_H)
    blank = prs.slide_layouts[6]

    meta = deck["meta"]
    cp = prs.core_properties
    cp.title = meta["title"]
    cp.subject = meta["title"]
    cp.author = "CKA Hands-On Training"
    cp.comments = "Generado con generador-diapositivas/"

    ctx_base = {"clase_n": meta["clase_n"], "curso": meta["curso"]}
    page = 0
    for s in deck["slides"]:
        page += 1
        slide = prs.slides.add_slide(blank)
        ctx = dict(ctx_base, page=page)
        RENDERERS[s["t"]](slide, s, ctx)

    prs.save(out_path)
    return prs, page
